pythonfrom pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

class Colors:
    """Palette de couleurs professionnelles"""
    
    @staticmethod
    def hex_to_rgb(hex_color):
        """Convertit une couleur hex en RGBColor"""
        if not hex_color:
            return RGBColor(255, 255, 255)
        
        hex_color = str(hex_color).lstrip('#')
        if len(hex_color) == 6:
            try:
                return RGBColor(
                    int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16)
                )
            except:
                return RGBColor(255, 255, 255)
        return RGBColor(255, 255, 255)
    
    @staticmethod
    def get_alignment(align_str):
        """Convertit string en PP_ALIGN"""
        if not align_str:
            return PP_ALIGN.LEFT
        
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        return align_map.get(str(align_str).lower(), PP_ALIGN.LEFT)
    
    @staticmethod
    def get_anchor(anchor_str):
        """Convertit string en MSO_ANCHOR"""
        if not anchor_str:
            return MSO_ANCHOR.TOP
        
        anchor_map = {
            'top': MSO_ANCHOR.TOP,
            'middle': MSO_ANCHOR.MIDDLE,
            'bottom': MSO_ANCHOR.BOTTOM
        }
        return anchor_map.get(str(anchor_str).lower(), MSO_ANCHOR.TOP)

class Formatter:
    """Utilitaires de formatage de texte"""
    
    @staticmethod
    def format_textbox(textbox, config, default_font="Arial"):
        """
        Applique le formatage à une textbox
        
        config: {
            'fontSize': int,
            'bold': bool,
            'color': str (hex),
            'align': str,
            'font': str,
            'anchor': str
        }
        """
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # Marges internes réduites
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        
        # Ancrage vertical
        anchor = config.get('anchor', 'top')
        text_frame.vertical_anchor = Colors.get_anchor(anchor)
        
        # Formatage du paragraphe
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = Colors.get_alignment(config.get('align', 'left'))
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(6)
            paragraph.line_spacing = 1.2
            
            # Formatage du run
            for run in paragraph.runs:
                run.font.name = config.get('font', default_font)
                run.font.size = Pt(config.get('fontSize', 16))
                run.font.bold = config.get('bold', False)
                
                if 'color' in config and config['color']:
                    run.font.color.rgb = Colors.hex_to_rgb(config['color'])
    
    @staticmethod
    def add_bullet_points(text_frame, items, config, default_font="Arial"):
        """Ajoute des bullet points formatés"""
        if not items or len(items) == 0:
            return
        
        text_frame.clear()
        text_frame.word_wrap = True
        
        # Marges
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        
        for i, item in enumerate(items):
            if not item:
                continue
            
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = str(item)
            p.level = 0
            
            # Bullet
            if config.get('bullet', True):
                p.bullet = True
            
            # Espacement
            p.space_before = Pt(3)
            p.space_after = Pt(3)
            p.line_spacing = 1.15
            p.alignment = Colors.get_alignment(config.get('align', 'left'))
            
            # Police
            for run in p.runs:
                run.font.name = config.get('font', default_font)
                run.font.size = Pt(config.get('fontSize', 16))
                run.font.bold = config.get('bold', False)
                
                if 'color' in config and config['color']:
                    run.font.color.rgb = Colors.hex_to_rgb(config['color'])
