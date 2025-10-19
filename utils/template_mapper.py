# utils/template_mapper.py
from pptx import Presentation
from pptx.util import Inches, Pt

class TemplateMapper:
    """Mappe les exercices aux slides du template et remplit les placeholders."""
    
    def __init__(self, template_path):
        self.template = Presentation(template_path)
        
        # Mapping type d'exercice → index de slide template
        self.slide_mapping = {
            'qcm': 2,                    # Slide 3 (index 2)
            'vrai_faux': 4,              # Slide 5 (index 4)
            'cas_pratique': 6,           # Slide 7 (index 6)
            'mise_en_situation': 8,      # Slide 9 (index 8)
            'exercice_pratique': 6,      # Slide 7 aussi
            'redaction': 2               # Slide 3
        }
    
    def find_text_shapes(self, slide):
        """Trouve toutes les zones de texte dans une slide."""
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_shapes.append(shape)
        return text_shapes
    
    def clear_and_fill_shape(self, shape, new_text):
        """Vide complètement un shape et le remplit avec le nouveau texte."""
        if not shape.has_text_frame:
            return
        
        text_frame = shape.text_frame
        
        # Vider tout le contenu
        text_frame.clear()
        
        # Ajouter le nouveau texte
        p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        p.text = new_text
        p.level = 0
    
    def fill_exercise_slide(self, slide, exercise_data):
        """Remplit une slide d'exercice avec les données."""
        text_shapes = self.find_text_shapes(slide)
        
        # Trier les shapes par position Y (de haut en bas)
        text_shapes.sort(key=lambda s: s.top)
        
        if len(text_shapes) < 3:
            print(f"⚠️ Pas assez de zones de texte trouvées : {len(text_shapes)}")
            return
        
        # Shape 0 : Généralement le numéro (1, 2, 3)
        # Shape 1 : Le TITRE de l'exercice
        # Shape 2+ : Consigne et contenu
        
        # Remplir le titre (shape 1 ou 2 selon la structure)
        titre = exercise_data.get('titre', 'Exercice')
        for i, shape in enumerate(text_shapes):
            if i == 1 or 'Rédiger' in shape.text or 'Créer' in shape.text or 'Exercice' in shape.text:
                self.clear_and_fill_shape(shape, titre)
                print(f"✅ Titre remplacé : {titre[:50]}...")
                break
        
        # Remplir la consigne
        consigne = exercise_data.get('consigne', '')
        if consigne:
            for shape in text_shapes:
                if 'Consigne' in shape.text:
                    self.clear_and_fill_shape(shape, f"Consigne : {consigne}")
                    print(f"✅ Consigne remplacée")
                    break
        
        # Gérer les types spécifiques
        ex_type = exercise_data.get('type', '')
        
        if ex_type == 'qcm':
            self._fill_qcm(slide, exercise_data, text_shapes)
        elif ex_type == 'vrai_faux':
            self._fill_vrai_faux(slide, exercise_data, text_shapes)
        else:
            # Pour exercice_pratique et mise_en_situation
            self._fill_generic(slide, exercise_data, text_shapes)
    
    def _fill_qcm(self, slide, data, text_shapes):
        """Remplit un QCM."""
        # Trouver la zone pour la question et les choix
        for shape in text_shapes:
            if 'Objectifs' in shape.text or 'Respecter' in shape.text or 'Maîtriser' in shape.text:
                tf = shape.text_frame
                tf.clear()
                
                # Question
                question = data.get('question', '')
                if question:
                    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                    p.text = question
                    p.level = 0
                    p.space_after = Pt(12)
                    
                    # Police en gras pour la question
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(16)
                
                # Choix
                for i, choix in enumerate(data.get('choix', [])):
                    p = tf.add_paragraph()
                    p.text = f"{chr(65+i)}. {choix}"
                    p.level = 0
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)
                    
                    for run in p.runs:
                        run.font.size = Pt(14)
                
                print(f"✅ QCM rempli avec {len(data.get('choix', []))} choix")
                break
    
    def _fill_vrai_faux(self, slide, data, text_shapes):
        """Remplit un Vrai/Faux."""
        for shape in text_shapes:
            if 'Objectifs' in shape.text or 'Respecter' in shape.text or 'Extraire' in shape.text:
                tf = shape.text_frame
                tf.clear()
                
                for aff in data.get('affirmations', []):
                    p = tf.paragraphs[0] if len(tf.paragraphs) == 1 else tf.add_paragraph()
                    p.text = f"• {aff.get('affirmation', '')}"
                    p.level = 0
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)
                    
                    for run in p.runs:
                        run.font.size = Pt(14)
                
                print(f"✅ Vrai/Faux rempli avec {len(data.get('affirmations', []))} affirmations")
                break
    
    def _fill_generic(self, slide, data, text_shapes):
        """Remplit les exercices génériques."""
        # Remplacer les objectifs pédagogiques par le contexte si disponible
        contexte = data.get('contexte', '')
        if contexte:
            for shape in text_shapes:
                if 'Objectifs' in shape.text or 'Respecter' in shape.text:
                    self.clear_and_fill_shape(shape, contexte)
                    print(f"✅ Contexte ajouté")
                    break