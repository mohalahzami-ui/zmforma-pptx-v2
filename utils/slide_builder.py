# utils/slide_builder.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile, os

class PresentationBuilder:
    def __init__(self, data, template_url=None):
        self.data = data
        self.slides_data = data.get('slides', [])
        self.default_font = data.get('theme', {}).get('font', 'Calibri')
        
        # Charger le template
        template_path = self._get_template_path(template_url)
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            print(f"✅ Template chargé : {template_path}")
            print(f"   Nombre de slides dans le template : {len(self.prs.slides)}")
            print(f"   Nombre de layouts disponibles : {len(self.prs.slide_layouts)}")
        else:
            print("⚠️ Template non trouvé, création présentation vide")
            self.prs = Presentation()
        
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
    
    def _get_template_path(self, template_url):
        """Retourne le chemin du template local."""
        local_path = os.path.join(
            os.path.dirname(__file__),
            "Formation.pptx"
        )
        return local_path if os.path.exists(local_path) else None
    
    def build(self):
        """Génère le PowerPoint complet."""
        # Supprimer toutes les slides sauf la première
        initial_slides = len(self.prs.slides)
        print(f"📊 Slides initiales : {initial_slides}")
        
        while len(self.prs.slides) > 1:
            rId = self.prs.slides._sldIdLst[1].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[1]
        
        print(f"📊 Après nettoyage : {len(self.prs.slides)} slide(s)")
        print(f"📊 Génération de {len(self.slides_data)} slides d'exercices...")
        
        # Générer les slides
        for i, slide_data in enumerate(self.slides_data):
            print(f"\n--- SLIDE {i+2} ---")
            print(f"Type : {slide_data.get('type')}")
            print(f"Titre : {slide_data.get('titre', '')[:50]}...")
            
            self._add_slide_from_scratch(slide_data)
            print(f"✅ Slide {i+2} créée")
        
        # Sauvegarder
        tmp_out = tempfile.mkstemp(suffix='.pptx')[1]
        self.prs.save(tmp_out)
        print(f"\n✅ PPTX généré avec {len(self.prs.slides)} slides : {tmp_out}")
        return tmp_out
    
    def _add_slide_from_scratch(self, slide_data):
        """Crée une nouvelle slide avec contenu."""
        # Utiliser le layout BLANK (généralement index 6)
        layout_idx = 6 if len(self.prs.slide_layouts) > 6 else 0
        layout = self.prs.slide_layouts[layout_idx]
        
        print(f"   Layout utilisé : index {layout_idx}")
        
        slide = self.prs.slides.add_slide(layout)
        
        # Données
        titre = slide_data.get('titre', 'Exercice')
        ex_type = slide_data.get('type', 'exercice_pratique')
        
        print(f"   Ajout du titre : {titre[:30]}...")
        
        # ===== TITRE =====
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.margin_top = Pt(10)
        title_frame.margin_bottom = Pt(10)
        
        p = title_frame.paragraphs[0]
        p.text = titre
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = self.default_font
        p.font.color.rgb = RGBColor(26, 26, 26)
        
        print(f"   ✓ Titre ajouté")
        
        # ===== CONTENU SELON TYPE =====
        if ex_type == 'qcm':
            self._add_qcm_content(slide, slide_data)
        elif ex_type == 'vrai_faux':
            self._add_vrai_faux_content(slide, slide_data)
        else:
            self._add_generic_content(slide, slide_data)
    
    def _add_qcm_content(self, slide, data):
        """Ajoute contenu QCM."""
        question = data.get('question', '')
        choix = data.get('choix', [])
        
        print(f"   QCM - Question : {question[:30] if question else 'VIDE'}...")
        print(f"   QCM - Nombre de choix : {len(choix)}")
        
        # QUESTION
        if question:
            q_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), Inches(9), Inches(1.2)
            )
            tf = q_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = question
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.name = self.default_font
            print(f"   ✓ Question ajoutée")
        
        # CHOIX
        if choix and len(choix) > 0:
            choices_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.2), Inches(9), Inches(2)
            )
            tf = choices_box.text_frame
            tf.word_wrap = True
            
            for i, choix_text in enumerate(choix):
                if i > 0:
                    p = tf.add_paragraph()
                else:
                    p = tf.paragraphs[0]
                
                p.text = f"{chr(65+i)}. {choix_text}"
                p.font.size = Pt(14)
                p.font.name = self.default_font
                p.space_before = Pt(6)
                p.space_after = Pt(6)
            
            print(f"   ✓ {len(choix)} choix ajoutés")
    
    def _add_vrai_faux_content(self, slide, data):
        """Ajoute contenu Vrai/Faux."""
        affirmations = data.get('affirmations', [])
        
        print(f"   Vrai/Faux - Nombre d'affirmations : {len(affirmations)}")
        
        if affirmations and len(affirmations) > 0:
            aff_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), Inches(9), Inches(3.5)
            )
            tf = aff_box.text_frame
            tf.word_wrap = True
            
            for i, aff in enumerate(affirmations):
                if i > 0:
                    p = tf.add_paragraph()
                else:
                    p = tf.paragraphs[0]
                
                aff_text = aff.get('affirmation', '') if isinstance(aff, dict) else str(aff)
                p.text = f"• {aff_text}"
                p.font.size = Pt(14)
                p.font.name = self.default_font
                p.space_before = Pt(8)
                p.space_after = Pt(8)
            
            print(f"   ✓ {len(affirmations)} affirmations ajoutées")
    
    def _add_generic_content(self, slide, data):
        """Ajoute contenu générique."""
        contexte = data.get('contexte', '')
        consigne = data.get('consigne', '')
        
        print(f"   Générique - Contexte : {contexte[:30] if contexte else 'VIDE'}...")
        print(f"   Générique - Consigne : {consigne[:30] if consigne else 'VIDE'}...")
        
        # CONTEXTE
        if contexte:
            ctx_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), Inches(9), Inches(1.5)
            )
            tf = ctx_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = contexte
            p.font.size = Pt(14)
            p.font.name = self.default_font
            print(f"   ✓ Contexte ajouté")
        
        # CONSIGNE
        if consigne:
            cons_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.5), Inches(9), Inches(1.5)
            )
            tf = cons_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"Consigne : {consigne}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.name = self.default_font
            print(f"   ✓ Consigne ajoutée")