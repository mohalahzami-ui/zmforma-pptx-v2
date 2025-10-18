from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

class Colors:
    """Palette de couleurs & helpers"""
    @staticmethod
    def hex_to_rgb(hex_color):
        if not hex_color:
            return RGBColor(255, 255, 255)
        hex_color = str(hex_color).lstrip('#')
        if len(hex_color) == 6:
            try:
                return RGBColor(
                    int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16)
                )
            except Exception:
                return RGBColor(255, 255, 255)
        return RGBColor(255, 255, 255)

    @staticmethod
    def get_alignment(align_str):
        if not align_str:
            return PP_ALIGN.LEFT
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        return align_map.get(str(align_str).lower(), PP_ALIGN.LEFT)

    @staticmethod
    def get_anchor(anchor_str):
        if not anchor_str:
            return MSO_ANCHOR.TOP
        anchor_map = {
            'top': MSO_ANCHOR.TOP,
            'middle': MSO_ANCHOR.MIDDLE,
            'bottom': MSO_ANCHOR.BOTTOM
        }
        return anchor_map.get(str(anchor_str).lower(), MSO_ANCHOR.TOP)


class Formatter:
    """Formatage texte + bullets (avec autosize)"""

    @staticmethod
    def _apply_run_style(run, config, default_font="Arial"):
        run.font.name = config.get('font', default_font)
        run.font.size = Pt(config.get('fontSize', 16))
        run.font.bold = config.get('bold', False)
        color = config.get('color')
        if color:
            run.font.color.rgb = Colors.hex_to_rgb(color)

    @staticmethod
    def format_textbox(textbox, config, default_font="Arial"):
        """
        Applique le formatage, active l'autofit du texte pour éviter le chevauchement.
        """
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Marges internes
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)

        # Ancrage vertical
        text_frame.vertical_anchor = Colors.get_anchor(config.get('anchor', 'top'))

        # Autofit : laisse PowerPoint réduire la taille si nécessaire
        try:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass  # Certaines versions ne supportent pas l'énum MSO_AUTO_SIZE

        # Paragraphe / alignement
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = Colors.get_alignment(config.get('align', 'left'))
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(4)
            paragraph.line_spacing = 1.18
            for run in paragraph.runs:
                Formatter._apply_run_style(run, config, default_font)

    @staticmethod
    def add_bullet_points(text_frame, items, config, default_font="Arial"):
        """Ajoute une liste à puces formatée (avec autosize)."""
        if not items:
            return

        text_frame.clear()
        text_frame.word_wrap = True

        text_frame.margin_left = Inches(0.10)
        text_frame.margin_right = Inches(0.10)
        text_frame.margin_top = Inches(0.06)
        text_frame.margin_bottom = Inches(0.06)

        try:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass

        for i, item in enumerate(items):
            if not item:
                continue
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = str(item)
            p.level = 0
            if config.get('bullet', True):
                p.bullet = True

            p.space_before = Pt(2)
            p.space_after = Pt(2)
            p.line_spacing = 1.15
            p.alignment = Colors.get_alignment(config.get('align', 'left'))

            for run in p.runs:
                Formatter._apply_run_style(run, config, default_font)
