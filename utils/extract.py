# utils/extract.py
from io import BytesIO
import fitz  # PyMuPDF
from pptx import Presentation

def extract_text_pdf(file_bytes: bytes) -> dict:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    texts = []
    for page in doc:
        texts.append(page.get_text("text"))
    result = {
        "kind": "pdf",
        "pages": len(doc),
        "text": "\n\n".join(texts).strip()
    }
    doc.close()
    return result

def extract_text_pptx(file_bytes: bytes) -> dict:
    prs = Presentation(BytesIO(file_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text)
    return {
        "kind": "pptx",
        "slides": len(prs.slides),
        "text": "\n\n".join(texts).strip()
    }
