pythonfrom pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import tempfile
import os
import requests
from io import BytesIO
from PIL import Image
from .styles import Colors, Formatter

class PresentationBuilder:
    """
    Constructeur de présentation PowerPoint professionnel
    Format 1920×1080 (16:9)
    """
    
    def __init__(self, data):
        self.data = data
        self.slides_data = data.get('slides', [])
        self.theme = data.get('theme', {})
        self.default_font = self.theme.get('font', 'Arial')
        
        # Initialiser la présentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        
    def build(self):
        """Construit la présentation complète"""
        print(f"🔨 Construction de {len(self.slides_data)} slides...")
        
        for i, slide_data in enumerate(self.slides_data):
            try:
                slide_type = slide_data.get('type', 'generic')
                print(f"  • Slide {i+1}/{len(self.slides_data)}: {slide_type}")
                
                # Dispatcher vers le bon builder
                if slide_type == 'cover':
                    self._build_cover(slide_data)
                elif slide_type == 'section':
                    self._build_section(slide_data)
                elif slide_type == 'qcm':
                    self._build_qcm(slide_data)
                elif slide_type == 'correction':
                    self._build_correction(slide_data)
                elif slide_type == 'vrai_faux':
                    self._build_vrai_faux(slide_data)
                elif slide_type == 'cas_pratique':
                    self._build_cas_pratique(slide_data)
                elif slide_type == 'mise_en_situation':
                    self._build_mise_en_situation(slide_data)
                elif slide_type == 'objectifs':
                    self._build_objectifs(slide_data)
                else:
                    self._build_generic(slide_data)
                    
            except Exception as e:
                print(f"❌ Erreur slide {i+1}: {str(e)}")
                import traceback
                print(traceback.format_exc())
                # Créer une slide d'erreur
                self._build_error_slide(i+1, str(e))
        
        # Sauvegarder
        temp_path = os.path.join(tempfile.gettempdir(), 'presentation_zmforma.pptx')
        self.prs.save(temp_path)
        print(f"✅ Présentation sauvegardée: {temp_path}")
        
        return temp_path
    
    def _build_cover(self, data):
        """Construit une slide de couverture"""
        layout = data.get('layout', {})
        bg_color = data.get('background', 'FFFFFF')
        
        # Créer slide vierge
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Fond coloré
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = Colors.hex_to_rgb(bg_color)
        
        # Titre
        if 'title' in layout:
            self._add_textbox(slide, layout['title'])
        
        # Sous-titre
        if 'subtitle' in layout:
            self._add_textbox(slide, layout['subtitle'])
        
        # Image
        if 'image' in layout and layout['image']:
            self._add_image(slide, layout['image'])
    
    def _build_section(self, data):
        """Construit une slide de section"""
        layout = data.get('layout', {})
        bg_color = data.get('background', 'F8FAFC')
        
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Fond
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = Colors.hex_to_rgb(bg_color)
        
        # Barre d'accent
        if 'accent_bar' in layout:
            self._add_shape(slide, layout['accent_bar'])
        
        # Kicker
        if 'kicker' in layout and layout['kicker']:
            self._add_textbox(slide, layout['kicker'])
        
        # Titre
        if 'title' in layout:
            self._add_textbox(slide, layout['title'])
        
        # Description
        if 'description' in layout and layout['description']:
            self._add_textbox(slide, layout['description'])
        
        # Image
        if 'image' in layout and layout['image']:
            self._add_image(slide, layout['image'])
    
    def _build_qcm(self, data):
        """Construit une slide QCM"""
        layout = data.get('layout', {})
        bg_color = data.get('background', 'FFFFFF')
        
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Fond
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = Colors.hex_to_rgb(bg_color)
        
        # Barre d'accent
        if 'accent_bar' in layout:
            self._add_shape(slide, layout['accent_bar'])
        
        # Kicker
        if 'kicker' in layout:
            self._add_textbox(slide, layout['kicker'])
        
        # Question
        if 'question' in layout:
            self._add_textbox(slide, layout['question'])
        
        # Contexte
        if 'context' in layout and layout['context']:
            self._add_textbox(slide, layout['context'])
        
        # Choix (bullets)
        if 'choices' in layout:
            self._add_bullets(slide, layout['choices'])
        
        # Image
        if 'image' in layout and layout['image']:
            self._add_image(slide, layout['image'])
    
    def _build_correction(self, data):
        """Construit une slide de correction"""
        layout = data.get('layout', {})
        bg_color = data.get('background', 'F0FDF4')
        
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Fond
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = Colors.hex_to_rgb(bg_color)
        
        # Barre
        if 'accent_bar' in layout:
            self._add_shape(slide, layout['accent_bar'])
        
        # Label
        if 'label' in layout:
            self._add_textbox(slide, layout['label'])
        
        # Titre/Réponse
        if 'answer' in layout:
            self._add_textbox(slide, layout['answer'])
        
        if 'answer_text' in layout:
            self._add_textbox(slide, layout['answer_text'])
        
        if 'title' in layout:
            self._add_textbox(slide, layout['title'])
        
        # Explication
        if 'explanation' in layout:
            self._add_textbox(slide, layout['explanation'])
        
        # Corrections (bullets)
        if 'corrections' in layout:
            self._add_bullets(slide, layout['corrections'])
        
        # Éléments (bullets)
        if 'elements' in layout:
            self._add_bullets(slide, layout['elements'])
    
    def _build_vrai_faux(self, data):
        """Slide vrai/faux (même structure que QCM)"""
        self._build_qcm(data)
    
    def _build_cas_pratique(self, data):
        """Slide cas pratique"""
        self._build_qcm(data)
    
    def _build_mise_en_situation(self, data):
        """Slide mise en situation"""
        self._build_qcm(data)
    
    def _build_objectifs(self, data):
        """Slide objectifs"""
        self._build_correction(data)
    
    def _build_generic(self, data):
        """Slide générique"""
        layout = data.get('layout', {})
        bg_color = data.get('background', 'FFFFFF')
        
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = Colors.hex_to_rgb(bg_color)
        
        # Ajouter tous les éléments du layout
        for key, element in layout.items():
            if element and isinstance(element, dict):
                if 'items' in element:
                    self._add_bullets(slide, element)
                elif 'text' in element:
                    self._add_textbox(slide, element)
                elif 'url' in element:
                    self._add_image(slide, element)
                elif 'fill' in element:
                    self._add_shape(slide, element)
    
    def _build_error_slide(self, slide_num, error_msg):
        """Créer une slide d'erreur"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 240, 240)
        
        # Titre erreur
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.8)
        )
        title_box.text = f"❌ Erreur - Slide {slide_num}"
        Formatter.format_textbox(title_box, {
            'fontSize': 32,
            'bold': True,
            'color': 'D32F2F',
            'align': 'center'
        })
        
        # Message d'erreur
        error_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9.0), Inches(2.0)
        )
        error_box.text = str(error_msg)[:500]
        Formatter.format_textbox(error_box, {
            'fontSize': 14,
            'color': '666666',
            'align': 'left'
        })
    
    # ==================== MÉTHODES UTILITAIRES ====================
    
    def _add_textbox(self, slide, config):
        """Ajoute une textbox formatée"""
        if not config or 'text' not in config:
            return
        
        x = Inches(config.get('x', 0.5))
        y = Inches(config.get('y', 1.0))
        w = Inches(config.get('w', 5.0))
        h = Inches(config.get('h', 1.0))
        
        textbox = slide.shapes.add_textbox(x, y, w, h)
        textbox.text = str(config['text'])
        
        Formatter.format_textbox(textbox, config, self.default_font)
    
    def _add_bullets(self, slide, config):
        """Ajoute une liste à puces"""
        if not config or 'items' not in config:
            return
        
        items = config['items']
        if not items or len(items) == 0:
            return
        
        x = Inches(config.get('x', 0.5))
        y = Inches(config.get('y', 1.0))
        w = Inches(config.get('w', 5.0))
        h = Inches(config.get('h', 2.0))
        
        textbox = slide.shapes.add_textbox(x, y, w, h)
        Formatter.add_bullet_points(textbox.text_frame, items, config, self.default_font)
    
    def _add_shape(self, slide, config):
        """Ajoute une forme (rectangle, etc.)"""
        if not config:
            return
        
        x = Inches(config.get('x', 0))
        y = Inches(config.get('y', 0))
        w = Inches(config.get('w', 1))
        h = Inches(config.get('h', 1))
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y, w, h
        )
        
        # Remplissage
        if 'fill' in config:
            shape.fill.solid()
            shape.fill.fore_color.rgb = Colors.hex_to_rgb(config['fill'])
        
        # Pas de bordure
        shape.line.fill.background()
    
    def _add_image(self, slide, config):
        """Ajoute une image (URL ou chemin local)"""
        if not config or 'url' not in config:
            return
        
        url = config['url']
        x = Inches(config.get('x', 5.0))
        y = Inches(config.get('y', 1.0))
        w = Inches(config.get('w', 3.0))
        h = Inches(config.get('h', 2.0))
        
        try:
            if url.startswith('http'):
                # Télécharger l'image
                response = requests.get(url, timeout=10)
                response.raise_for_status()
                image_stream = BytesIO(response.content)
                
                # Vérifier que c'est une image valide
                img = Image.open(image_stream)
                img.verify()
                
                # Réouvrir pour l'ajout
                image_stream.seek(0)
                slide.shapes.add_picture(image_stream, x, y, width=w, height=h)
            else:
                # Chemin local
                if os.path.exists(url):
                    slide.shapes.add_picture(url, x, y, width=w, height=h)
        except Exception as e:
            print(f"⚠️ Impossible d'ajouter l'image {url}: {str(e)}")
